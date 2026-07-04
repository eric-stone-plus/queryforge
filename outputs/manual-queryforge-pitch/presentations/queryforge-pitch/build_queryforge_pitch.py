#!/usr/bin/env python3
"""Build the QueryForge pitch deck with python-pptx."""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path("/Users/ericstone/Downloads/data-agent")
WORKSPACE = ROOT / "outputs/manual-queryforge-pitch/presentations/queryforge-pitch"
ASSET_DIR = WORKSPACE / "assets"
SCREENSHOT = ASSET_DIR / "live-app.png"
OUT_PATH = ROOT / "assets/QueryForge-Pitch.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb(hex_value: str) -> RGBColor:
    h = hex_value.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


BG = rgb("#1B2A4A")
CARD = rgb("#1F3056")
CARD_DARK = rgb("#17243F")
AMBER = rgb("#FFB300")
WHITE = rgb("#FFFFFF")
MUTED = rgb("#B0BEC5")
MUTED_2 = rgb("#7F90A8")
GREEN = rgb("#4CAF50")
RED = rgb("#FF6B6B")
BLUE = rgb("#58A6FF")
LINE = rgb("#2D426D")
INK = rgb("#0F172A")

FONT = "Aptos"
FONT_DISPLAY = "Aptos Display"


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_rect(slide, x, y, w, h, fill=CARD, line=None, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=18,
    color=WHITE,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    all_caps=False,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = valign

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.08
    run = p.add_run()
    run.text = text.upper() if all_caps else text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraphs(slide, x, y, w, h, lines, size=16, color=MUTED, gap=5):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, item in enumerate(lines):
        text, item_color, bold, item_size = (
            item if isinstance(item, tuple) else (item, color, False, size)
        )
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        p.line_spacing = 1.12
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(item_size)
        run.font.bold = bold
        run.font.color.rgb = item_color
    return box


def add_title(slide, kicker, title, support=None):
    add_text(
        slide,
        Inches(0.72),
        Inches(0.42),
        Inches(3.8),
        Inches(0.26),
        kicker,
        size=10,
        color=AMBER,
        bold=True,
        all_caps=True,
    )
    add_text(
        slide,
        Inches(0.72),
        Inches(0.72),
        Inches(9.6),
        Inches(0.72),
        title,
        size=32,
        color=WHITE,
        bold=True,
        font=FONT_DISPLAY,
    )
    add_rect(slide, Inches(0.72), Inches(1.45), Inches(2.15), Pt(3), fill=AMBER)
    if support:
        add_text(
            slide,
            Inches(0.72),
            Inches(1.62),
            Inches(10.6),
            Inches(0.38),
            support,
            size=13,
            color=MUTED,
        )


def add_footer(slide, page):
    add_text(
        slide,
        Inches(0.72),
        Inches(7.12),
        Inches(5.5),
        Inches(0.18),
        "QueryForge | ClawHunt Builder Camp 2026",
        size=8.5,
        color=MUTED_2,
    )
    add_text(
        slide,
        Inches(12.0),
        Inches(7.12),
        Inches(0.65),
        Inches(0.18),
        f"{page}/9",
        size=8.5,
        color=MUTED_2,
        align=PP_ALIGN.RIGHT,
    )


def add_card(slide, x, y, w, h, title, body, accent=AMBER, title_size=16):
    add_rect(slide, x, y, w, h, fill=CARD, line=LINE)
    add_rect(slide, x, y, Pt(4), h, fill=accent)
    add_text(slide, x + Inches(0.22), y + Inches(0.20), w - Inches(0.44), Inches(0.28), title, size=title_size, color=WHITE, bold=True)
    if body:
        add_paragraphs(slide, x + Inches(0.22), y + Inches(0.63), w - Inches(0.44), h - Inches(0.78), body, size=13.2, color=MUTED, gap=5)


def add_metric_card(slide, x, y, w, h, label, value, note, accent=AMBER):
    add_rect(slide, x, y, w, h, fill=CARD, line=LINE)
    add_text(slide, x + Inches(0.18), y + Inches(0.18), w - Inches(0.36), Inches(0.22), label, size=10.5, color=MUTED, bold=True, all_caps=True)
    add_text(slide, x + Inches(0.18), y + Inches(0.50), w - Inches(0.36), Inches(0.46), value, size=27, color=accent, bold=True, font=FONT_DISPLAY)
    add_text(slide, x + Inches(0.18), y + Inches(1.06), w - Inches(0.36), Inches(0.35), note, size=10.5, color=MUTED)


def add_flow_step(slide, x, y, w, h, step, title, note, accent=AMBER):
    add_rect(slide, x, y, w, h, fill=CARD, line=LINE)
    add_text(slide, x + Inches(0.18), y + Inches(0.14), Inches(0.6), Inches(0.28), step, size=11, color=accent, bold=True)
    add_text(slide, x + Inches(0.18), y + Inches(0.50), w - Inches(0.36), Inches(0.30), title, size=15.5, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.18), y + Inches(0.92), w - Inches(0.36), h - Inches(1.05), note, size=11.5, color=MUTED)


def add_arrow_text(slide, x, y):
    add_text(slide, x, y, Inches(0.34), Inches(0.18), "->", size=14, color=AMBER, bold=True, align=PP_ALIGN.CENTER)


def add_image_frame(slide, image_path, x, y, w, h, line=LINE):
    add_rect(slide, x - Pt(2), y - Pt(2), w + Pt(4), h + Pt(4), fill=rgb("#F7FAFC"), line=line)
    slide.shapes.add_picture(str(image_path), x, y, width=w, height=h)


def prepare_screenshot_crops():
    if not SCREENSHOT.exists():
        return {}
    im = Image.open(SCREENSHOT).convert("RGB")
    crops = {
        "kpis": (0, 58, 1338, 156),
        "input": (54, 936, 832, 988),
        "dashboard": (885, 155, 1348, 1000),
        "library": (1346, 0, 1600, 295),
    }
    out = {}
    for name, box in crops.items():
        target = ASSET_DIR / f"crop-{name}.png"
        im.crop(box).save(target)
        out[name] = target
    return out


def slide_01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    add_text(slide, Inches(0.72), Inches(0.82), Inches(2.2), Inches(0.24), "PITCH DECK", size=10, color=AMBER, bold=True, all_caps=True)
    add_text(slide, Inches(0.72), Inches(1.32), Inches(7.6), Inches(0.85), "QueryForge", size=58, color=AMBER, bold=True, font=FONT_DISPLAY)
    add_rect(slide, Inches(0.72), Inches(2.35), Inches(2.55), Pt(4), fill=AMBER)
    add_text(
        slide,
        Inches(0.72),
        Inches(2.76),
        Inches(7.8),
        Inches(0.82),
        "Natural-language business data queries with charts on demand.",
        size=25,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        Inches(0.72),
        Inches(3.78),
        Inches(7.2),
        Inches(0.72),
        "Business users ask a question, QueryForge fetches the data, and the answer appears as KPI cards, charts, or reusable dashboard views.",
        size=15,
        color=MUTED,
    )

    add_rect(slide, Inches(8.95), Inches(1.10), Inches(3.45), Inches(3.8), fill=CARD, line=LINE)
    add_text(slide, Inches(9.26), Inches(1.42), Inches(2.5), Inches(0.25), "LIVE CONTEXT", size=10, color=AMBER, bold=True, all_caps=True)
    context = [
        ("Built in 72 hours", WHITE, True, 17),
        ("ClawHunt Builder Camp 2026 hackathon", MUTED, False, 12.5),
        ("", MUTED, False, 5),
        ("Powered by MiMo v2.5 Pro", WHITE, True, 17),
        ("Model used in the live demo", MUTED, False, 12.5),
        ("", MUTED, False, 5),
        ("Demo data", WHITE, True, 17),
        ("10K orders | 500 products | 1,000 users | 8 regions | 20 categories", MUTED, False, 12.5),
    ]
    add_paragraphs(slide, Inches(9.26), Inches(1.82), Inches(2.78), Inches(2.72), context, size=12.5, gap=3)

    add_text(slide, Inches(0.72), Inches(6.44), Inches(7.5), Inches(0.24), "queryforge-production-8d6f.up.railway.app", size=12.5, color=BLUE, bold=True)
    add_footer(slide, 1)
    return slide


def slide_02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(
        slide,
        "PAIN POINT",
        "The bottleneck is request handling, not analysis.",
        "When business teams need answers, every clarification becomes another handoff."
    )

    add_card(
        slide,
        Inches(0.72),
        Inches(2.24),
        Inches(5.75),
        Inches(3.92),
        "Analyst queue",
        [
            "Same metric is requested by multiple teams in slightly different words.",
            "Questions arrive as ad hoc requests, follow-ups, and revised definitions.",
            "Analysts spend attention on recreating views instead of interpreting the business."
        ],
        accent=RED,
        title_size=18,
    )
    add_card(
        slide,
        Inches(6.86),
        Inches(2.24),
        Inches(5.75),
        Inches(3.92),
        "Business waiting loop",
        [
            "Revenue, category, channel, and customer questions come up during operating reviews.",
            "Teams need a chart they can inspect immediately, not another request thread.",
            "Metric definitions have to stay consistent as more people reuse the answer."
        ],
        accent=BLUE,
        title_size=18,
    )
    add_text(slide, Inches(0.72), Inches(6.55), Inches(10.9), Inches(0.28), "The product opportunity: make common business questions self-service while keeping analyst-defined metrics reusable.", size=14, color=WHITE, bold=True)
    add_footer(slide, 2)
    return slide


def slide_03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(
        slide,
        "SOLUTION",
        "Analysts define the metric library once; business users self-serve.",
        "QueryForge turns the analyst role from repeated report builder into metric owner."
    )

    x0, y0 = Inches(0.72), Inches(2.20)
    w, h, gap = Inches(3.68), Inches(3.70), Inches(0.33)
    add_flow_step(
        slide,
        x0,
        y0,
        w,
        h,
        "01",
        "Analyst setup",
        "Define preset metrics once: revenue trend, category margin, top products, channel orders, and repeat buyer views.",
        AMBER,
    )
    add_arrow_text(slide, x0 + w + Inches(0.02), y0 + Inches(1.78))
    add_flow_step(
        slide,
        x0 + w + gap,
        y0,
        w,
        h,
        "02",
        "Business questions",
        "Users type a natural-language question or rerun a saved metric from the library.",
        BLUE,
    )
    add_arrow_text(slide, x0 + 2 * w + gap + Inches(0.05), y0 + Inches(1.78))
    add_flow_step(
        slide,
        x0 + 2 * (w + gap),
        y0,
        w,
        h,
        "03",
        "Reusable answers",
        "Charts, KPI cards, progress streaming, and self-correction keep the answer usable in front of the business.",
        GREEN,
    )

    add_text(slide, Inches(0.72), Inches(6.28), Inches(11.3), Inches(0.30), "Core capability: plain-language data access that preserves analyst-defined business meaning.", size=14, color=WHITE, bold=True)
    add_footer(slide, 3)
    return slide


def slide_04(prs, crops):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(
        slide,
        "PRODUCT DEMO",
        "The live product already combines questions, KPI cards, and dashboard panels.",
        "Screenshots are from the deployed QueryForge demo with actual demo data."
    )

    bullets = [
        ("Natural-language entry", WHITE, True, 15),
        ("A user types a business question and gets a charted answer.", MUTED, False, 12.2),
        ("", MUTED, False, 4),
        ("8 KPI cards", WHITE, True, 15),
        ("Revenue, AOV, margin, repeat rate, completion, refund, basket, active buyers.", MUTED, False, 12.2),
        ("", MUTED, False, 4),
        ("6 dashboard panels", WHITE, True, 15),
        ("Region bar, category pie, monthly trend, channel bar, top products, user segments.", MUTED, False, 12.2),
        ("", MUTED, False, 4),
        ("Preset metric library", WHITE, True, 15),
        ("Analysts create reusable metrics that business users can rerun.", MUTED, False, 12.2),
    ]
    add_paragraphs(slide, Inches(0.72), Inches(2.08), Inches(3.55), Inches(4.6), bullets, size=12.2, gap=3)

    if crops:
        add_image_frame(slide, crops["kpis"], Inches(4.70), Inches(1.92), Inches(7.85), Inches(0.58))
        add_image_frame(slide, crops["input"], Inches(4.70), Inches(2.78), Inches(4.62), Inches(0.31))
        add_image_frame(slide, crops["library"], Inches(9.62), Inches(2.78), Inches(2.93), Inches(1.12))
        add_image_frame(slide, crops["dashboard"], Inches(4.70), Inches(3.34), Inches(4.18), Inches(3.36))
        add_text(slide, Inches(9.12), Inches(4.18), Inches(3.1), Inches(0.40), "Dashboard remains visible while the user asks or reruns questions.", size=13, color=WHITE, bold=True)
        add_text(slide, Inches(9.12), Inches(4.80), Inches(3.1), Inches(0.90), "The product view is built around business application: operating health, regional performance, channel mix, product focus, and user segmentation.", size=12.3, color=MUTED)
    else:
        add_text(slide, Inches(4.70), Inches(2.8), Inches(6.8), Inches(1), "Live screenshot unavailable in this environment.", size=20, color=RED, bold=True)

    add_footer(slide, 4)
    return slide


def slide_05(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(
        slide,
        "DATA PROOF",
        "Eight KPI cards are backed by the demo dataset.",
        "The dashboard summarizes the operating view business users can reuse."
    )

    metrics = [
        ("Total revenue", "\u00a5232M", "30-month demo history", AMBER),
        ("AOV", "\u00a523K", "Average order value", AMBER),
        ("Margin", "46.7%", "All-category average", GREEN),
        ("Repeat rate", "100%", "1,000 active buyers", GREEN),
        ("Completion", "66.5%", "6,650 completed orders", BLUE),
        ("Refund", "16.6%", "1,664 refunded orders", RED),
        ("Basket", "2.5 items", "Average items per order", AMBER),
        ("Active buyers", "1,000", "Across 20 categories", BLUE),
    ]
    x0, y0 = Inches(0.72), Inches(2.08)
    w, h = Inches(2.83), Inches(1.50)
    gx, gy = Inches(0.25), Inches(0.34)
    for i, (label, value, note, accent) in enumerate(metrics):
        col = i % 4
        row = i // 4
        add_metric_card(slide, x0 + col * (w + gx), y0 + row * (h + gy), w, h, label, value, note, accent)

    add_rect(slide, Inches(0.72), Inches(5.82), Inches(11.88), Inches(0.62), fill=CARD_DARK, line=LINE)
    add_text(slide, Inches(0.96), Inches(6.02), Inches(11.4), Inches(0.22), "Demo data scope: 10K orders | 500 products | 1,000 users | 8 regions | 20 categories", size=14, color=WHITE, bold=True)
    add_footer(slide, 5)
    return slide


def slide_06(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(
        slide,
        "HOW IT WORKS",
        "Question to answer, with progress visible throughout.",
        "The flow is designed for business users: ask, watch progress, inspect the chart, reuse the metric."
    )

    y = Inches(2.12)
    x = Inches(0.72)
    w = Inches(2.25)
    h = Inches(1.50)
    steps = [
        ("1", "Ask", "Type a business question in natural language.", AMBER),
        ("2", "Stream", "Progress updates show what is happening.", BLUE),
        ("3", "Answer", "Data appears as charts, tables, or KPI cards.", GREEN),
        ("4", "Reuse", "Save useful answers into the metric library.", AMBER),
    ]
    for i, (num, title, note, accent) in enumerate(steps):
        xx = x + i * Inches(3.03)
        add_flow_step(slide, xx, y, w, h, num, title, note, accent)
        if i < len(steps) - 1:
            add_arrow_text(slide, xx + w + Inches(0.36), y + Inches(0.63))

    add_rect(slide, Inches(2.45), Inches(4.58), Inches(8.35), Inches(1.18), fill=CARD_DARK, line=LINE)
    add_text(slide, Inches(2.76), Inches(4.84), Inches(2.05), Inches(0.30), "Self-correction loop", size=15, color=WHITE, bold=True)
    add_text(slide, Inches(5.05), Inches(4.82), Inches(5.25), Inches(0.45), "If a query fails, QueryForge fixes the attempt and retries so the user can stay focused on the business question.", size=12.5, color=MUTED)
    add_text(slide, Inches(0.72), Inches(6.30), Inches(11.4), Inches(0.32), "No handoff is required for routine data questions once the metrics have been set up.", size=14, color=WHITE, bold=True)
    add_footer(slide, 6)
    return slide


def slide_07(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(
        slide,
        "USE CASES",
        "The demo scenarios map directly to operating questions.",
        "Each scenario starts with a business question and ends with a reusable visual answer."
    )

    cases = [
        ("Regional performance", "How are regions contributing to revenue over time?", "Output: region distribution and monthly trend."),
        ("Category margin", "Which categories carry the strongest margin?", "Output: category comparison view."),
        ("Product focus", "Which products are driving the most revenue?", "Output: top products table."),
        ("Customer behavior", "Which buyers and segments show repeat activity?", "Output: repeat buyer and user segment views."),
    ]
    x0, y0 = Inches(0.72), Inches(2.12)
    w, h = Inches(5.78), Inches(1.38)
    gx, gy = Inches(0.34), Inches(0.32)
    for i, (title, q, out) in enumerate(cases):
        col = i % 2
        row = i // 2
        x = x0 + col * (w + gx)
        y = y0 + row * (h + gy)
        add_rect(slide, x, y, w, h, fill=CARD, line=LINE)
        add_text(slide, x + Inches(0.22), y + Inches(0.18), w - Inches(0.44), Inches(0.26), title, size=16, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.22), y + Inches(0.58), w - Inches(0.44), Inches(0.28), q, size=12.5, color=MUTED)
        add_text(slide, x + Inches(0.22), y + Inches(0.96), w - Inches(0.44), Inches(0.24), out, size=12.0, color=AMBER, bold=True)

    add_text(slide, Inches(0.72), Inches(6.34), Inches(11.6), Inches(0.28), "The same surface supports quick exploration and repeated management reporting.", size=14, color=WHITE, bold=True)
    add_footer(slide, 7)
    return slide


def slide_08(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(
        slide,
        "WHAT'S NEXT",
        "Move from a hackathon demo to a business-ready tool.",
        "The roadmap extends the existing product without adding unsupported market or revenue claims."
    )

    lanes = [
        ("Deployment", "Harden the live environment, improve onboarding, and make the demo easier to hand off to evaluators or early users."),
        ("More data sources", "Connect additional business datasets beyond the ecommerce demo while preserving the same question-to-chart experience."),
        ("Team features", "Add shared metric libraries, saved views, and collaboration controls for analysts and business stakeholders."),
    ]
    x0, y0 = Inches(0.72), Inches(2.20)
    for i, (title, body) in enumerate(lanes):
        y = y0 + i * Inches(1.38)
        add_rect(slide, x0, y, Inches(11.9), Inches(1.04), fill=CARD, line=LINE)
        add_text(slide, x0 + Inches(0.28), y + Inches(0.22), Inches(2.0), Inches(0.28), title, size=17, color=AMBER, bold=True)
        add_text(slide, x0 + Inches(2.55), y + Inches(0.20), Inches(8.65), Inches(0.46), body, size=13.0, color=MUTED)

    add_footer(slide, 8)
    return slide


def slide_09(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    add_text(slide, Inches(0.72), Inches(0.88), Inches(2.8), Inches(0.24), "CTA", size=10, color=AMBER, bold=True, all_caps=True)
    add_text(slide, Inches(0.72), Inches(1.42), Inches(7.2), Inches(0.84), "QueryForge", size=55, color=AMBER, bold=True, font=FONT_DISPLAY)
    add_rect(slide, Inches(0.72), Inches(2.42), Inches(2.55), Pt(4), fill=AMBER)
    add_text(slide, Inches(0.72), Inches(2.86), Inches(8.8), Inches(0.64), "Let business teams ask for data directly while analysts own the reusable metrics.", size=23, color=WHITE, bold=True)

    add_rect(slide, Inches(0.72), Inches(4.28), Inches(11.4), Inches(1.40), fill=CARD, line=LINE)
    add_text(slide, Inches(1.02), Inches(4.58), Inches(1.55), Inches(0.24), "Live demo", size=12, color=MUTED, bold=True, all_caps=True)
    add_text(slide, Inches(2.58), Inches(4.54), Inches(8.6), Inches(0.30), "queryforge-production-8d6f.up.railway.app", size=18, color=BLUE, bold=True)
    add_text(slide, Inches(1.02), Inches(5.14), Inches(1.55), Inches(0.24), "GitHub", size=12, color=MUTED, bold=True, all_caps=True)
    add_text(slide, Inches(2.58), Inches(5.10), Inches(8.6), Inches(0.30), "github.com/eric-stone-plus/queryforge", size=18, color=BLUE, bold=True)

    add_text(slide, Inches(0.72), Inches(6.42), Inches(8.2), Inches(0.24), "Built at ClawHunt Builder Camp 2026 in 72 hours | MiMo v2.5 Pro", size=12, color=MUTED)
    add_footer(slide, 9)
    return slide


def build():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    crops = prepare_screenshot_crops()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "QueryForge Pitch"
    prs.core_properties.subject = "Business self-service data query tool"
    prs.core_properties.author = "QueryForge"

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs, crops)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved {OUT_PATH}")


if __name__ == "__main__":
    build()
