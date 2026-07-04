#!/usr/bin/env python3
"""Regenerate the Chinese QueryForge pitch deck with python-pptx."""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path("/Users/ericstone/Downloads/data-agent")
OUT_PATH = ROOT / "assets/QueryForge-Pitch.pptx"
SOURCE_ASSET_DIR = ROOT / "outputs/manual-queryforge-pitch/presentations/queryforge-pitch/assets"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


BG = rgb("#1a1f2e")
ACCENT = rgb("#d4a853")
SUBTITLE = rgb("#9a9a9a")
CARD = rgb("#2a2f3e")
TEXT = rgb("#f4f1e9")
TEXT_SOFT = rgb("#c8c4bc")
LINE = rgb("#3b4356")
RED = rgb("#d97171")
GREEN = rgb("#64b47c")
BLUE = rgb("#80aee8")

FONT = "Microsoft YaHei"
FONT_ALT = "PingFang SC"


def set_ea_font(run, typeface=FONT):
    run.font.name = typeface
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = rpr.find(qn(tag))
        if node is None:
            node = rpr.makeelement(qn(tag), {})
            rpr.append(node)
        node.set("typeface", typeface)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=18,
    color=TEXT,
    bold=False,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.06
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    set_ea_font(run)
    return box


def add_paragraphs(slide, x, y, w, h, items, size=14, color=TEXT_SOFT, gap=7):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)

    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, item_color, item_bold, item_size = item
        else:
            text, item_color, item_bold, item_size = item, color, False, size
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.12
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(item_size)
        run.font.bold = item_bold
        run.font.color.rgb = item_color
        set_ea_font(run)
    return box


def add_rect(slide, x, y, w, h, fill=CARD, line=LINE, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_rule(slide, x, y, w, color=ACCENT):
    add_rect(slide, x, y, w, Pt(2), fill=color, line=None)


def add_title(slide, section, title, subtitle=None):
    add_text(slide, Inches(0.72), Inches(0.42), Inches(3.6), Inches(0.24), section, 10.5, ACCENT, True)
    add_text(slide, Inches(0.72), Inches(0.76), Inches(10.2), Inches(0.54), title, 30, TEXT, True)
    add_rule(slide, Inches(0.72), Inches(1.42), Inches(2.18))
    if subtitle:
        add_text(slide, Inches(0.72), Inches(1.60), Inches(10.8), Inches(0.42), subtitle, 13, SUBTITLE)


def add_footer(slide, text):
    add_text(slide, Inches(0.72), Inches(7.06), Inches(9.6), Inches(0.18), text, 8.5, SUBTITLE)


def add_card(slide, x, y, w, h, title, body_lines, accent=ACCENT):
    add_rect(slide, x, y, w, h, fill=CARD, line=LINE)
    add_rect(slide, x, y, Pt(4), h, fill=accent, line=None)
    add_text(slide, x + Inches(0.22), y + Inches(0.20), w - Inches(0.44), Inches(0.28), title, 15.5, TEXT, True)
    add_paragraphs(
        slide,
        x + Inches(0.22),
        y + Inches(0.63),
        w - Inches(0.44),
        h - Inches(0.78),
        body_lines,
        size=12.6,
        color=TEXT_SOFT,
        gap=5,
    )


def add_metric(slide, x, y, w, h, label, value, note, accent=ACCENT):
    add_rect(slide, x, y, w, h, fill=CARD, line=LINE)
    add_text(slide, x + Inches(0.18), y + Inches(0.16), w - Inches(0.36), Inches(0.22), label, 10.8, SUBTITLE, True)
    add_text(slide, x + Inches(0.18), y + Inches(0.46), w - Inches(0.36), Inches(0.46), value, 25.5, accent, True)
    add_text(slide, x + Inches(0.18), y + Inches(1.05), w - Inches(0.36), Inches(0.28), note, 10.5, TEXT_SOFT)


def add_image(slide, path, x, y, w, h):
    add_rect(slide, x - Pt(2), y - Pt(2), w + Pt(4), h + Pt(4), fill=rgb("#f6f7f9"), line=LINE)
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def crop_assets():
    crops = {
        "kpis": SOURCE_ASSET_DIR / "crop-kpis.png",
        "input": SOURCE_ASSET_DIR / "crop-input.png",
    }
    dashboard = SOURCE_ASSET_DIR / "crop-dashboard.png"
    if dashboard.exists():
        image = Image.open(dashboard).convert("RGB")
        generated = {
            "region": (0, 136, 463, 365),
            "category": (0, 384, 226, 620),
            "channel": (236, 384, 463, 620),
        }
        for name, box in generated.items():
            target = SOURCE_ASSET_DIR / f"ppt-{name}.png"
            image.crop(box).save(target)
            crops[name] = target
    return crops


def slide_01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    add_text(slide, Inches(0.72), Inches(1.18), Inches(7.3), Inches(0.88), "QueryForge", 58, ACCENT, True)
    add_rule(slide, Inches(0.72), Inches(2.20), Inches(2.72))
    add_text(slide, Inches(0.72), Inches(2.62), Inches(7.4), Inches(0.52), "分析师定义一次，业务自助用数据", 25, TEXT, True)
    add_text(
        slide,
        Inches(0.72),
        Inches(3.30),
        Inches(7.0),
        Inches(0.78),
        "把重复取数、可视化、异常归因和行动建议放进一个业务工作流，让团队不用为每个经营问题重新排队。",
        15,
        SUBTITLE,
    )

    add_rect(slide, Inches(8.36), Inches(1.34), Inches(3.92), Inches(3.22), fill=CARD, line=LINE)
    proof = [
        ("10,000", "订单验证", ACCENT),
        ("8", "地区覆盖", BLUE),
        ("20", "品类覆盖", GREEN),
    ]
    for i, (value, label, color) in enumerate(proof):
        y = Inches(1.72 + i * 0.82)
        add_text(slide, Inches(8.70), y, Inches(1.52), Inches(0.36), value, 24, color, True)
        add_text(slide, Inches(10.25), y + Inches(0.08), Inches(1.55), Inches(0.24), label, 12.5, TEXT_SOFT)
    add_text(slide, Inches(8.70), Inches(4.10), Inches(2.7), Inches(0.25), "电商经营数据实测", 12, SUBTITLE)

    add_footer(slide, "定位：面向业务团队的自助数据分析工作台")


def slide_02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "02 痛点", "业务要数据，总要等", "不是没有数据，而是每个业务问题都要通过分析师手工中转。")

    scenarios = [
        ("运营晨会", ["想知道华东销售额下滑在哪个品类。", "需要补维度、等日报、再追问原因。"], RED),
        ("市场复盘", ["要同时看渠道收入、客单价和退款率。", "同一批数据被重新组合好几次。"], ACCENT),
        ("管理追问", ["经营指标低于预期，需要当天解释。", "分析师临时加急，常规分析被打断。"], BLUE),
    ]
    x = Inches(0.72)
    for title, body, color in scenarios:
        add_card(slide, x, Inches(2.28), Inches(3.66), Inches(2.10), title, body, color)
        x += Inches(4.04)

    add_rect(slide, Inches(0.72), Inches(5.15), Inches(11.82), Inches(0.72), fill=rgb("#222838"), line=LINE)
    add_text(slide, Inches(0.98), Inches(5.36), Inches(2.35), Inches(0.24), "典型等待链路", 13, ACCENT, True)
    add_text(slide, Inches(3.58), Inches(5.36), Inches(7.8), Inches(0.24), "提需求 → 排期 → 拉数 → 改口径 → 再解释", 14, TEXT, True)
    add_text(slide, Inches(0.72), Inches(6.35), Inches(10.8), Inches(0.24), "每次重复请求，都会占用分析师一次上下文切换，也推迟业务决策。", 13.5, TEXT_SOFT)
    add_footer(slide, "痛点：业务需要的是即时经营判断，分析师不应该反复做同类取数。")


def slide_03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "03 方案", "分析师定义一次，业务自助使用", "QueryForge 把分析师的领域知识沉淀为可复用的指标资产。")

    add_card(
        slide,
        Inches(0.72),
        Inches(2.24),
        Inches(3.35),
        Inches(2.75),
        "分析师",
        ["定义指标口径。", "维护常用经营维度。", "把高频问题沉淀成可复用指标。"],
        BLUE,
    )
    add_text(slide, Inches(4.30), Inches(3.22), Inches(0.55), Inches(0.35), "→", 23, ACCENT, True)
    add_card(
        slide,
        Inches(4.92),
        Inches(2.24),
        Inches(3.15),
        Inches(2.75),
        "可信指标库",
        ["地区月度销售额趋势。", "品类利润率对比。", "渠道营收对比。", "复购用户排行。"],
        ACCENT,
    )
    add_text(slide, Inches(8.30), Inches(3.22), Inches(0.55), Inches(0.35), "→", 23, ACCENT, True)
    add_card(
        slide,
        Inches(8.92),
        Inches(2.24),
        Inches(3.35),
        Inches(2.75),
        "业务团队",
        ["用自然语言提问。", "直接得到图表和指标。", "继续追问原因和下一步动作。"],
        GREEN,
    )
    add_text(slide, Inches(0.72), Inches(5.72), Inches(10.9), Inches(0.28), "结果：口径由分析师把关，使用由业务自助完成。", 16, TEXT, True)
    add_footer(slide, "从“提需求等三天”到“自己问自己看”。")


def slide_04(prs, crops):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "04 产品能力", "取数 → 可视化 → 归因 → 建议，一个流程走完", "业务用户从问题出发，得到可复用的经营答案。")

    stages = [
        ("取数", "示例问题：各地区月度销售额趋势", "自动匹配地区、月份和收入指标。", ACCENT),
        ("可视化", "结果视图：地区营收分布、渠道订单分布", "用图表直接呈现经营结构。", BLUE),
        ("归因", "追问问题：退款率为什么升高？", "按地区、品类和渠道定位异常来源。", RED),
        ("建议", "行动建议：优先处理高退款品类和低完成率渠道", "把洞察转成下一步运营动作。", GREEN),
    ]
    x0 = Inches(0.72)
    y0 = Inches(2.10)
    w = Inches(2.83)
    for i, (stage, example, note, color) in enumerate(stages):
        x = x0 + i * Inches(3.05)
        add_rect(slide, x, y0, w, Inches(2.02), fill=CARD, line=LINE)
        add_text(slide, x + Inches(0.18), y0 + Inches(0.16), w - Inches(0.36), Inches(0.30), stage, 16, color, True)
        add_text(slide, x + Inches(0.18), y0 + Inches(0.62), w - Inches(0.36), Inches(0.48), example, 11.8, TEXT, True)
        add_text(slide, x + Inches(0.18), y0 + Inches(1.24), w - Inches(0.36), Inches(0.45), note, 10.8, TEXT_SOFT)

    if crops.get("input"):
        add_image(slide, crops["input"], Inches(0.84), Inches(4.80), Inches(5.65), Inches(0.38))
        add_text(slide, Inches(0.84), Inches(5.33), Inches(5.8), Inches(0.24), "产品截图：业务人员直接用自然语言输入分析需求", 11, SUBTITLE)
    if crops.get("region"):
        add_image(slide, crops["region"], Inches(7.05), Inches(4.45), Inches(2.62), Inches(1.28))
    if crops.get("channel"):
        add_image(slide, crops["channel"], Inches(9.94), Inches(4.45), Inches(1.94), Inches(1.28))
    add_footer(slide, "产品能力：从经营问题到可执行动作，减少重复沟通。")


def slide_05(prs, crops):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "05 真实验证", "10K 订单、8 地区、20 品类，撑起 8 个核心指标", "演示数据覆盖订单、商品、用户、地区、品类和渠道。")

    scope = [("10,000", "订单"), ("8", "地区"), ("20", "品类")]
    for i, (value, label) in enumerate(scope):
        x = Inches(0.72 + i * 2.12)
        add_text(slide, x, Inches(2.06), Inches(1.65), Inches(0.42), value, 28, ACCENT, True)
        add_text(slide, x, Inches(2.58), Inches(1.45), Inches(0.22), label, 11.5, TEXT_SOFT)

    metrics = [
        ("总营收", "¥23,256万", "30 个月累计", ACCENT),
        ("客单价", "¥23,256", "平均每单", ACCENT),
        ("毛利率", "46.7%", "全品类均值", GREEN),
        ("复购率", "100%", "人均 10 单", GREEN),
        ("完成率", "66.5%", "6,650 单完成", BLUE),
        ("退款率", "16.6%", "1,664 单退款", RED),
        ("连带率", "2.5件", "每单平均", ACCENT),
        ("活跃买家", "1,000", "覆盖 20 品类", BLUE),
    ]
    x0, y0 = Inches(0.72), Inches(3.18)
    w, h = Inches(2.82), Inches(1.18)
    for i, (label, value, note, color) in enumerate(metrics):
        col = i % 4
        row = i // 4
        add_metric(slide, x0 + col * Inches(3.06), y0 + row * Inches(1.36), w, h, label, value, note, color)

    if crops.get("kpis"):
        add_image(slide, crops["kpis"], Inches(0.80), Inches(6.12), Inches(11.25), Inches(0.82))
    add_footer(slide, "验证重点：指标不是静态展示，而是可被业务继续追问和复用。")


def slide_06(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "06 双重价值", "分析师更专注，业务团队更独立", "同一套指标口径服务两类用户，减少重复劳动，也加快经营决策。")

    add_card(
        slide,
        Inches(0.72),
        Inches(2.06),
        Inches(5.52),
        Inches(3.82),
        "对分析师",
        [
            "从重复拉数转向指标治理。",
            "把常用问题沉淀为可复用资产。",
            "减少口径解释和反复修改。",
            "把时间留给归因、策略和深度洞察。",
        ],
        BLUE,
    )
    add_card(
        slide,
        Inches(6.78),
        Inches(2.06),
        Inches(5.52),
        Inches(3.82),
        "对业务团队",
        [
            "运营、市场、财务随时自助查看。",
            "同一页面完成提问、图表和指标复用。",
            "追问速度更快，决策不过夜。",
            "用统一口径讨论经营结果。",
        ],
        GREEN,
    )
    add_text(slide, Inches(0.72), Inches(6.36), Inches(11.1), Inches(0.28), "关键变化：分析师不再是每个问题的手工中转，而是全公司可信指标的定义者。", 14, TEXT, True)
    add_footer(slide, "双重价值：效率提升来自同一套可信业务语义。")


def slide_07(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "07 应用场景", "高频经营问题，都可以变成自助分析流程", "从日常运营到管理复盘，每个角色都能用同一套指标快速追问。")

    cases = [
        ("电商运营", "今日华东收入是否低于上月同期？", "定位地区、品类和商品表现，调整库存与促销。", ACCENT),
        ("市场团队", "不同渠道的订单和收入贡献如何变化？", "复盘活动效果，把预算投向更有效渠道。", BLUE),
        ("管理层", "本周总营收、毛利率和完成率是否健康？", "形成经营日报，快速发现异常并追问原因。", GREEN),
        ("财务", "收入、成本和利润分布有什么结构变化？", "支持预算执行、成本控制和利润分析。", RED),
    ]
    for i, (role, question, action, color) in enumerate(cases):
        col = i % 2
        row = i // 2
        x = Inches(0.72 + col * 6.06)
        y = Inches(2.10 + row * 2.02)
        add_card(slide, x, y, Inches(5.52), Inches(1.66), role, [question, action], color)
    add_text(slide, Inches(0.72), Inches(6.42), Inches(10.8), Inches(0.24), "适合对象：有数据团队、但业务侧高频提需求的中型企业。", 13.5, TEXT_SOFT)
    add_footer(slide, "应用场景：电商运营、市场团队、管理层、财务都能直接使用。")


def slide_08(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    add_text(slide, Inches(0.72), Inches(1.18), Inches(7.3), Inches(0.86), "QueryForge", 56, ACCENT, True)
    add_rule(slide, Inches(0.72), Inches(2.20), Inches(2.72))
    add_text(slide, Inches(0.72), Inches(2.66), Inches(8.5), Inches(0.78), "让数据分析师做更有价值的事\n让业务团队自己找到答案", 24, TEXT, True)
    add_text(slide, Inches(0.72), Inches(4.08), Inches(8.2), Inches(0.38), "产品演示", 13, SUBTITLE, True)
    add_text(slide, Inches(0.72), Inches(4.48), Inches(9.6), Inches(0.32), "queryforge-production-8d6f.up.railway.app", 17, ACCENT, True)
    add_text(slide, Inches(0.72), Inches(5.14), Inches(8.2), Inches(0.38), "代码仓库（GitHub）", 13, SUBTITLE, True)
    add_text(slide, Inches(0.72), Inches(5.54), Inches(9.6), Inches(0.32), "github.com/eric-stone-plus/queryforge", 16, TEXT_SOFT, True)

    add_rect(slide, Inches(9.15), Inches(1.50), Inches(2.98), Inches(3.60), fill=CARD, line=LINE)
    add_text(slide, Inches(9.46), Inches(1.84), Inches(2.1), Inches(0.32), "下一步", 17, ACCENT, True)
    add_paragraphs(
        slide,
        Inches(9.46),
        Inches(2.38),
        Inches(2.18),
        Inches(1.86),
        ["邀请真实业务团队试用。", "沉淀更多指标模板。", "验证日常经营复盘流程。"],
        size=12.4,
        color=TEXT_SOFT,
        gap=7,
    )
    add_footer(slide, "结束：从重复取数走向业务自助分析。")


def build():
    crops = crop_assets()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "QueryForge 中文路演稿"
    prs.core_properties.subject = "业务自助数据分析"
    prs.core_properties.author = "QueryForge"

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs, crops)
    slide_05(prs, crops)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved {OUT_PATH} with {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
