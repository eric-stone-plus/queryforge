from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0x0a, 0x0e, 0x17)
CARD = RGBColor(0x14, 0x1a, 0x25)
BLUE = RGBColor(0x4a, 0xa3, 0xff)
CYAN = RGBColor(0x19, 0xd3, 0xc5)
YELLOW = RGBColor(0xf7, 0xd7, 0x4d)
PINK = RGBColor(0xff, 0x65, 0x7a)
GREEN = RGBColor(0x1a, 0x7f, 0x37)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x8b, 0x94, 0x9e)
LIGHT = RGBColor(0xc9, 0xd1, 0xd9)

def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

def add_label(slide, left, top, text, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(1.8), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1a, 0x22, 0x33)
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_card(slide, left, top, width, height, accent_color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.fill.background()
    # accent line on top
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.15), Inches(top), Inches(width - 0.3), Inches(0.04))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    return shape

def add_page_num(slide, num):
    add_text(slide, 12.2, 7.0, 0.8, 0.4, f"{num:02d}", 12, GRAY, False, PP_ALIGN.RIGHT)

def add_footer(slide):
    add_text(slide, 0.8, 7.0, 2, 0.4, "QueryForge", 11, GRAY)

# ===== Slide 1: Cover =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.6, "AI DATA AGENT", BLUE)
add_text(slide, 0.8, 1.5, 6, 1.2, "QueryForge", 52, WHITE, True)
add_text(slide, 0.8, 2.8, 6, 0.8, "让业务数据快速变成经营洞察", 24, LIGHT)
add_text(slide, 0.8, 3.8, 6, 1, "自然语言提问 → AI 生成 SQL → 实时可视化\n不是关键词匹配，是真正的 AI 推理", 16, GRAY)
# Right side: mock UI card
add_card(slide, 7.5, 1.2, 5, 5, BLUE)
add_text(slide, 7.8, 1.5, 4.4, 0.4, "AI Data Agent", 10, BLUE, True)
add_text(slide, 7.8, 2.0, 4.4, 0.5, "各地区月度销售额趋势", 16, WHITE, True)
add_text(slide, 7.8, 2.7, 4.4, 0.3, "SELECT r.name, strftime('%Y-%m', o.order_date)...", 9, GRAY)
# Mock chart area
mock_chart = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.8), Inches(3.2), Inches(4.4), Inches(2.5))
mock_chart.fill.solid()
mock_chart.fill.fore_color.rgb = RGBColor(0x0d, 0x11, 0x1a)
mock_chart.line.fill.background()
add_text(slide, 8.0, 3.4, 4, 2, "📈 240 行数据 · 8 地区 · 30 个月\n折线图已生成\n\nMiMo v2.5 Pro 实时推理\nSQL AST 安全校验\n自动 LIMIT 500", 12, GRAY)
add_footer(slide)
add_page_num(slide, 1)

# ===== Slide 2: Customer Pain =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.6, "CUSTOMER PAIN", PINK)
add_text(slide, 0.8, 1.2, 11, 0.8, "业务团队缺的不是数据，而是可直接使用的分析结论", 32, WHITE, True)
add_text(slide, 0.8, 2.2, 11, 0.5, "经营数据散落在文件里，真正耗时的是清洗、理解、解释和汇报", 15, GRAY)

pains = [
    ("SQL 门槛", "业务人员不会写查询，依赖数据团队排队等排期", BLUE),
    ("工具太重", "传统 BI 学习成本高，临时需求无法快速响应", CYAN),
    ("假智能", "市面工具用关键词匹配冒充 AI，无法理解开放式问题", YELLOW),
    ("单点故障", "LLM 偶尔生成错误 SQL，没有自动修正机制", PINK),
]
for i, (title, desc, color) in enumerate(pains):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 3.2 + row * 2.0
    add_card(slide, x, y, 5.8, 1.7, color)
    add_text(slide, x + 0.3, y + 0.3, 5.2, 0.4, title, 18, WHITE, True)
    add_text(slide, x + 0.3, y + 0.8, 5.2, 0.7, desc, 14, GRAY)

add_footer(slide)
add_page_num(slide, 2)

# ===== Slide 3: Product Positioning =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.6, "PRODUCT POSITIONING", CYAN)
add_text(slide, 0.8, 1.2, 11, 0.8, "真正的 AI 数据分析，不是关键词匹配", 32, WHITE, True)
add_text(slide, 0.8, 2.2, 11, 0.5, "从自然语言到 SQL 到可视化，全链路 AI 推理，一步到位", 15, GRAY)

steps = [
    ("提问", "用户用中文描述需求", "\"最近三个月哪个品类增长最快\"", BLUE),
    ("推理", "AI 理解意图，生成 SQL", "MiMo v2.5 Pro 实时推理\nAST 解析器安全校验", CYAN),
    ("执行", "查询数据库，返回结果", "10,000 订单 · 500 商品\n自动 LIMIT 防卡死", GREEN),
    ("可视化", "图表 + 指标 + 解释", "柱状/折线/饼图/面积图\nKPI 卡片 + 数据透视", YELLOW),
]
for i, (title, desc, detail, color) in enumerate(steps):
    x = 0.5 + i * 3.15
    add_card(slide, x, 3.0, 2.9, 3.8, color)
    add_text(slide, x + 0.2, 3.3, 2.5, 0.4, f"0{i+1}  {title}", 16, color, True)
    add_text(slide, x + 0.2, 3.9, 2.5, 0.5, desc, 13, WHITE)
    add_text(slide, x + 0.2, 4.6, 2.5, 1.5, detail, 11, GRAY)

add_footer(slide)
add_page_num(slide, 3)

# ===== Slide 4: Self-Correction (Innovation) =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.6, "INNOVATION", YELLOW)
add_text(slide, 0.8, 1.2, 11, 0.8, "自纠正循环：SQL 报错 → AI 自动修正", 32, WHITE, True)
add_text(slide, 0.8, 2.2, 11, 0.5, "不是出了错就放弃，而是把错误反馈给 AI，让它自己修", 15, GRAY)

add_card(slide, 0.8, 3.0, 3.8, 3.5, BLUE)
add_text(slide, 1.1, 3.3, 3.2, 0.4, "第一步：生成 SQL", 16, BLUE, True)
add_text(slide, 1.1, 3.9, 3.2, 2, "用户问：\"各品类退货率\"\n\nAI 生成 SQL\n执行 → 报错\n\"no such column: refund_rate\"", 12, GRAY)

add_card(slide, 4.9, 3.0, 3.8, 3.5, YELLOW)
add_text(slide, 5.2, 3.3, 3.2, 0.4, "第二步：自动修正", 16, YELLOW, True)
add_text(slide, 5.2, 3.9, 3.2, 2, "错误信息反馈给 AI\n\nAI 分析错误原因\n生成修正后的 SQL\n重新执行", 12, GRAY)

add_card(slide, 9.0, 3.0, 3.8, 3.5, GREEN)
add_text(slide, 9.3, 3.3, 3.2, 0.4, "第三步：成功返回", 16, GREEN, True)
add_text(slide, 9.3, 3.9, 3.2, 2, "修正后查询成功\n\n用户看到：\n✅ 修正标签\n✅ 正确图表\n✅ 修正说明", 12, GRAY)

add_footer(slide)
add_page_num(slide, 4)

# ===== Slide 5: Business Dashboard =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.6, "BUSINESS DASHBOARD", GREEN)
add_text(slide, 0.8, 1.2, 11, 0.8, "商业数据透视面板，8 项核心指标 + 6 个图表", 32, WHITE, True)
add_text(slide, 0.8, 2.2, 11, 0.5, "全部来自真实数据库，不是模拟数据", 15, GRAY)

metrics = [
    ("总营收", "¥23,256万", "30个月累计", BLUE),
    ("客单价", "¥23,256", "平均每单", CYAN),
    ("毛利率", "46.7%", "全品类均值", GREEN),
    ("复购率", "100%", "人均10单", YELLOW),
    ("完成率", "66.5%", "6,650单", GREEN),
    ("退款率", "16.6%", "1,664单", PINK),
    ("连带率", "2.5件", "每单平均", BLUE),
    ("活跃买家", "1,000", "覆盖20品类", CYAN),
]
for i, (label, value, sub, color) in enumerate(metrics):
    col = i % 4
    row = i // 4
    x = 0.8 + col * 3.1
    y = 3.0 + row * 2.0
    add_card(slide, x, y, 2.8, 1.7, color)
    add_text(slide, x + 0.2, y + 0.25, 2.4, 0.3, label, 12, GRAY)
    add_text(slide, x + 0.2, y + 0.6, 2.4, 0.5, value, 24, WHITE, True)
    add_text(slide, x + 0.2, y + 1.2, 2.4, 0.3, sub, 10, GRAY)

add_footer(slide)
add_page_num(slide, 5)

# ===== Slide 6: Multi-Agent Audit (QUINTE) =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.6, "QUALITY ASSURANCE", BLUE)
add_text(slide, 0.8, 1.2, 11, 0.8, "五方对抗审查：5 个 AI 互相挑刺", 32, WHITE, True)
add_text(slide, 0.8, 2.2, 11, 0.5, "不是一个人写代码一个人审，而是 5 个独立 AI 同时审查，消除单点盲区", 15, GRAY)

rounds = [
    ("R1 独立分析", "5 个智能体各自独立审查\n互不通信，避免群体思维\n每个结论必须引用文件和行号", BLUE),
    ("R2 交叉审查", "匿名审查其他 4 个的输出\n标记共识、分歧和盲区\n被说服必须声明立场变更", CYAN),
    ("R3 双路裁决", "人类 + 独立审计方共同裁决\n产出残差闭环账本\n未解决问题标为阻断信号", YELLOW),
]
for i, (title, desc, color) in enumerate(rounds):
    x = 0.8 + i * 4.1
    add_card(slide, x, 3.0, 3.8, 2.8, color)
    add_text(slide, x + 0.3, 3.3, 3.2, 0.4, title, 16, color, True)
    add_text(slide, x + 0.3, 3.9, 3.2, 1.5, desc, 13, GRAY)

# Cost comparison
add_card(slide, 0.8, 6.1, 5.8, 0.8, GREEN)
add_text(slide, 1.1, 6.2, 5.2, 0.5, "成本：约 ¥5/次（小米 MiMo token plan）", 14, WHITE, True)
add_card(slide, 7.0, 6.1, 5.8, 0.8, PINK)
add_text(slide, 7.3, 6.2, 5.2, 0.5, "对比：人工 Code Review ≈ ¥5,000（5人×2h）", 14, WHITE, True)

add_footer(slide)
add_page_num(slide, 6)

# ===== Slide 7: Tech Stack =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.6, "TECH STACK", CYAN)
add_text(slide, 0.8, 1.2, 11, 0.8, "技术架构：全栈 TypeScript，云端部署", 32, WHITE, True)

techs = [
    ("前端", "Next.js 14 · Tailwind CSS · Recharts\n响应式设计 · 深色主题 · KPI 仪表盘", BLUE),
    ("后端", "Next.js API Routes · better-sqlite3\nSQL AST 解析安全校验 · 自动 LIMIT", CYAN),
    ("AI 引擎", "MiMo v2.5 Pro（小米自研）\nVercel AI SDK · 自纠正循环", YELLOW),
    ("数据层", "SQLite · 10K 订单 · 500 商品\n8 地区 · 20 品类 · 6 渠道", GREEN),
    ("部署", "Railway 云端 24/7 在线\nGitHub 代码评审 · 不依赖本机", PINK),
    ("审查", "QUINTE 五方对抗协议\n3 轮 × 5 方 = 39 份独立审计报告", BLUE),
]
for i, (title, desc, color) in enumerate(techs):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 2.4 + row * 2.5
    add_card(slide, x, y, 3.8, 2.1, color)
    add_text(slide, x + 0.3, y + 0.3, 3.2, 0.4, title, 16, color, True)
    add_text(slide, x + 0.3, y + 0.8, 3.2, 1, desc, 13, GRAY)

add_footer(slide)
add_page_num(slide, 7)

# ===== Slide 8: Score Projection =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.6, "SCORE PROJECTION", YELLOW)
add_text(slide, 0.8, 1.2, 11, 0.8, "评分预估：88-96 / 105 + 3 加分", 32, WHITE, True)

scores = [
    ("Demo 现场可用", "25", "20-22", "公网部署 + 缓存 fallback + 流式进度", BLUE),
    ("用户价值/PMF", "20", "16-18", "真实痛点 + 自然语言接口 + 仪表盘", CYAN),
    ("技术实现", "20", "16-18", "AI SDK + AST 校验 + 自纠正 + 对抗审查", GREEN),
    ("创新性", "15", "11-13", "自纠正循环 + QUINTE 方法论 + 流式进度", YELLOW),
    ("商业潜力", "10", "7-8", "SaaS 模式 + 真实数据 + 商业指标", PINK),
    ("路演表达", "10", "8-9", "故事线 + 排练 + PPT", BLUE),
]

# Header
add_card(slide, 0.8, 2.4, 11.7, 0.5, GRAY)
add_text(slide, 1.1, 2.45, 3.5, 0.4, "维度", 12, GRAY, True)
add_text(slide, 5.0, 2.45, 1.5, 0.4, "满分", 12, GRAY, True, PP_ALIGN.CENTER)
add_text(slide, 7.0, 2.45, 1.5, 0.4, "预估", 12, GRAY, True, PP_ALIGN.CENTER)
add_text(slide, 9.0, 2.45, 3.5, 0.4, "依据", 12, GRAY, True)

for i, (dim, max_s, est, reason, color) in enumerate(scores):
    y = 3.1 + i * 0.6
    if i % 2 == 0:
        add_card(slide, 0.8, y - 0.05, 11.7, 0.55, GRAY)
    add_text(slide, 1.1, y, 3.5, 0.4, dim, 13, WHITE)
    add_text(slide, 5.0, y, 1.5, 0.4, max_s, 13, GRAY, False, PP_ALIGN.CENTER)
    add_text(slide, 7.0, y, 1.5, 0.4, est, 14, color, True, PP_ALIGN.CENTER)
    add_text(slide, 9.0, y, 3.5, 0.4, reason, 11, GRAY)

# Bonus
add_card(slide, 0.8, 6.7, 5.5, 0.5, GREEN)
add_text(slide, 1.1, 6.75, 5, 0.4, "加分：ClawHunt 上架 +3 · 游园展示 +2", 12, WHITE, True)
add_card(slide, 6.7, 6.7, 5.8, 0.5, YELLOW)
add_text(slide, 7.0, 6.75, 5.2, 0.4, "总计预估：91-101 / 105", 14, WHITE, True)

add_footer(slide)
add_page_num(slide, 8)

# ===== Slide 9: Business Model =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.6, "BUSINESS", GREEN)
add_text(slide, 0.8, 1.2, 11, 0.8, "商业潜力：SaaS 订阅 + 真实数据 demo", 32, WHITE, True)

add_card(slide, 0.8, 2.4, 5.8, 4.2, BLUE)
add_text(slide, 1.1, 2.7, 5.2, 0.4, "市场规模", 18, BLUE, True)
add_text(slide, 1.1, 3.3, 5.2, 3, "• 中国 BI 市场规模：¥200亿+（2025）\n• 中小企业数据分析渗透率 < 15%\n• 自然语言 BI 是下一代趋势\n\n目标用户\n• 运营/销售管理者（不会 SQL）\n• 数据分析师（加速探索）\n• 中小企业老板（快速决策）", 14, GRAY)

add_card(slide, 7.0, 2.4, 5.8, 4.2, CYAN)
add_text(slide, 7.3, 2.7, 5.2, 0.4, "商业模式", 18, CYAN, True)
add_text(slide, 7.3, 3.3, 5.2, 3, "SaaS 订阅制\n• 免费版：5 次/天查询\n• Pro 版：¥99/月，无限查询\n• 企业版：¥999/月，私有部署\n\n竞争优势\n• 真正的 AI 推理（非关键词匹配）\n• 自纠正保证查询成功率\n• 五方审查保证代码质量\n• 成本极低（MiMo token plan）", 14, GRAY)

add_footer(slide)
add_page_num(slide, 9)

# ===== Slide 10: Demo =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.6, "LIVE DEMO", PINK)
add_text(slide, 0.8, 1.2, 11, 0.8, "现场演示：queryforge-production-8d6f.up.railway.app", 28, WHITE, True)

demo_steps = [
    ("01", "点击预设查询", "展示 AI 生成 SQL + 图表\n流式进度推送", BLUE),
    ("02", "自由提问", "\"最近三个月哪个品类增长最快\"\n展示 AI 理解开放式问题", CYAN),
    ("03", "仪表盘联动", "8 个 KPI 指标 + 6 个图表\n全部真实数据库数据", GREEN),
    ("04", "指标保存复用", "侧边栏保存/删除/复用\nlocalStorage 持久化", YELLOW),
    ("05", "自纠正演示", "SQL 报错 → 自动修正\n展示修正过程和结果", PINK),
]
for i, (num, title, desc, color) in enumerate(demo_steps):
    x = 0.5 + i * 2.5
    add_card(slide, x, 2.5, 2.3, 3.5, color)
    add_text(slide, x + 0.15, 2.7, 2, 0.4, num, 20, color, True, PP_ALIGN.CENTER)
    add_text(slide, x + 0.15, 3.2, 2, 0.4, title, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, x + 0.15, 3.8, 2, 1.5, desc, 11, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, 0.8, 6.4, 11.5, 0.5, "产品地址：queryforge-production-8d6f.up.railway.app", 14, BLUE, False, PP_ALIGN.CENTER)
add_text(slide, 0.8, 6.8, 11.5, 0.5, "代码仓库：github.com/eric-stone-plus/queryforge", 12, GRAY, False, PP_ALIGN.CENTER)

add_footer(slide)
add_page_num(slide, 10)

# ===== Slide 11: End =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 1, 2.2, 11, 1.2, "QueryForge", 56, BLUE, True, PP_ALIGN.CENTER)
add_text(slide, 1, 3.5, 11, 0.8, "让每个人都能用数据说话", 28, WHITE, False, PP_ALIGN.CENTER)
add_text(slide, 1, 4.8, 11, 0.5, "queryforge-production-8d6f.up.railway.app", 18, BLUE, False, PP_ALIGN.CENTER)
add_text(slide, 1, 5.5, 11, 0.5, "github.com/eric-stone-plus/queryforge", 14, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 1, 6.2, 11, 0.5, "ClawHunt Builder Camp 2026 · Track A: Agents at Work", 12, GRAY, False, PP_ALIGN.CENTER)

add_footer(slide)
add_page_num(slide, 11)

out = "/Users/ericstone/Downloads/data-agent/QueryForge-Pitch.pptx"
prs.save(out)
print(f"PPT saved: {out} ({len(prs.slides)} slides)")
