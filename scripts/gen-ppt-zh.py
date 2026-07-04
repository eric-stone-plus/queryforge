from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BG = RGBColor(0x1a, 0x1f, 0x2e)
TEXT = RGBColor(0xf0, 0xee, 0xe8)
SUBTEXT = RGBColor(0x9a, 0x9a, 0x9a)
ACCENT = RGBColor(0xd4, 0xa8, 0x53)
CARD_BG = RGBColor(0x2a, 0x2f, 0x3e)

TITLE_SIZE = Pt(30)
BODY_SIZE = Pt(15)
BIG_NUM_SIZE = Pt(48)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(1.0)

def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=BODY_SIZE, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': 'Microsoft YaHei'})
    rPr.append(ea)
    return txBox

def add_title(slide, text, top=MARGIN):
    return add_text(slide, MARGIN, top, Inches(10), Inches(0.6), text, size=TITLE_SIZE, bold=True)

def add_bullet_list(slide, items, left=MARGIN, top=Inches(2.0), width=Inches(10)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(len(items) * 0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"  —  {item}"
        run.font.size = BODY_SIZE
        run.font.color.rgb = TEXT
        run.font.name = "Arial"
        rPr = run._r.get_or_add_rPr()
        ea = rPr.makeelement(qn('a:ea'), {'typeface': 'Microsoft YaHei'})
        rPr.append(ea)
    return txBox

def add_big_number(slide, left, top, number, label):
    add_text(slide, left, top, Inches(3), Inches(0.8), str(number), size=BIG_NUM_SIZE, color=ACCENT, bold=True)
    add_text(slide, left, top + Inches(0.9), Inches(3), Inches(0.4), label, size=BODY_SIZE, color=SUBTEXT)

def add_card(slide, left, top, width, height, title, content):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    add_text(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.4), title, size=Pt(16), bold=True, color=ACCENT)
    add_text(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), height - Inches(1), content, size=BODY_SIZE, color=TEXT)

def add_bottom_text(slide, text):
    add_text(slide, MARGIN, Inches(6.5), Inches(10), Inches(0.4), text, size=Pt(12), color=SUBTEXT)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Page 1: 封面
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_text(s, MARGIN, Inches(2.5), Inches(10), Inches(1), "QueryForge", size=Pt(48), bold=True, color=TEXT)
add_text(s, MARGIN, Inches(3.5), Inches(10), Inches(0.5), "让业务团队自助取数，解放数据分析师", size=Pt(20), color=SUBTEXT)
add_bottom_text(s, "ClawHunt Builder Camp 2026 · 72 小时构建")

# Page 2: 痛点
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_title(s, "业务要数据，总要等")
add_bullet_list(s, [
    '运营问"华东上个月卖了多少" → 找分析师',
    '市场问"哪个品类增长最快" → 排期等报表',
    '老板问"为什么这个月下滑" → 加急再加急',
])
add_bottom_text(s, "数据分析师每周 60% 时间在重复拉数")

# Page 3: 方案
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_title(s, "分析师定义一次，业务自助使用")
add_bullet_list(s, [
    "分析师定义数据口径和指标",
    "业务用自然语言提问",
    "系统返回数据、图表、洞察",
])
add_bottom_text(s, '从"提需求等三天"到"自己问自己看"')

# Page 4: 产品能力
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_title(s, "从提问到决策，一个流程走完")
add_bullet_list(s, [
    '取数 — "各地区本月销售额是多少"',
    "可视化 — 自动生成图表，一键切换维度",
    '归因 — "为什么东区下降了" → 自动找原因',
    "建议 — 基于数据给出下一步行动",
])

# Page 5: 真实验证
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_title(s, "在真实数据上验证")
add_big_number(s, MARGIN, Inches(2.0), "10,000", "订单")
add_big_number(s, Inches(4.5), Inches(2.0), "8", "地区")
add_big_number(s, Inches(8), Inches(2.0), "20", "品类")
add_bullet_list(s, [
    "覆盖 8 个 KPI 指标",
    "支持跨地区、跨品类对比",
    "异常检测与趋势分析",
], top=Inches(4.5))

# Page 6: 双重价值
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_title(s, "一个产品，两方受益")
add_card(s, MARGIN, Inches(2.0), Inches(5), Inches(4), "分析师", "减少重复取数工作\n专注于指标定义和数据治理\n从'拉数的'变成'定标准的'")
add_card(s, Inches(7), Inches(2.0), Inches(5), Inches(4), "业务团队", "随时提问，不用排队\n自己看数据，自己做决策\n从'等报表'到'看数据'")

# Page 7: 应用场景
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_title(s, "谁需要这个工具")
add_bullet_list(s, [
    "电商运营 — 每天看销售、库存、转化",
    "市场团队 — 活动效果、渠道对比、用户画像",
    "管理层 — 经营日报、异常预警、决策支持",
    "财务 — 收入分析、成本归因、预算执行",
])

# Page 8: 结束
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_text(s, MARGIN, Inches(2.5), Inches(10), Inches(1), "QueryForge", size=Pt(48), bold=True, color=TEXT)
add_text(s, MARGIN, Inches(3.5), Inches(10), Inches(0.8), "让数据分析师做更有价值的事\n让业务团队自己找到答案", size=Pt(18), color=SUBTEXT)
add_bottom_text(s, "queryforge-production-8d6f.up.railway.app · github.com/eric-stone-plus/queryforge")

out = "/Users/ericstone/Downloads/data-agent/assets/QueryForge-Pitch.pptx"
prs.save(out)
print(f"PPT saved: {out} ({len(prs.slides)} slides)")
