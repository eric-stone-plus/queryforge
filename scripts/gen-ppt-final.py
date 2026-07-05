import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BG = RGBColor(0x1a, 0x1f, 0x2e)
TEXT = RGBColor(0xf0, 0xee, 0xe8)
SUBTEXT = RGBColor(0x9a, 0x9a, 0x9a)
ACCENT = RGBColor(0xd4, 0xa8, 0x53)
CARD = RGBColor(0x2a, 0x2f, 0x3e)
WHITE = RGBColor(0xff, 0xff, 0xff)

W = Inches(13.333)
H = Inches(7.5)
M = Inches(1.0)

def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

def txt(s, l, t, w, h, text, sz=Pt(15), c=TEXT, b=False, a=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.alignment = a
    r = p.add_run()
    r.text = text
    r.font.size = sz
    r.font.bold = b
    r.font.color.rgb = c
    r.font.name = "Arial"
    rp = r._r.get_or_add_rPr()
    ea = rp.makeelement(qn('a:ea'), {'typeface': 'Microsoft YaHei'})
    rp.append(ea)
    return box

def title(s, text, t=M):
    return txt(s, M, t, Inches(10), Inches(0.6), text, Pt(30), TEXT, True)

def bullets(s, items, t=Inches(2.0)):
    box = s.shapes.add_textbox(M, t, Inches(10), Inches(len(items)*0.55))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"  —  {item}"
        r.font.size = Pt(15)
        r.font.color.rgb = TEXT
        r.font.name = "Arial"
        rp = r._r.get_or_add_rPr()
        ea = rp.makeelement(qn('a:ea'), {'typeface': 'Microsoft YaHei'})
        rp.append(ea)

def big(s, l, t, num, label):
    txt(s, l, t, Inches(3), Inches(0.8), str(num), Pt(48), ACCENT, True)
    txt(s, l, t+Inches(0.9), Inches(3), Inches(0.4), label, Pt(15), SUBTEXT)

def card(s, l, t, w, h, title, body):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = CARD
    sh.line.fill.background()
    txt(s, l+Inches(0.3), t+Inches(0.2), w-Inches(0.6), Inches(0.4), title, Pt(16), ACCENT, True)
    txt(s, l+Inches(0.3), t+Inches(0.7), w-Inches(0.6), h-Inches(1), body, Pt(14), TEXT)

def footer(s, text):
    txt(s, M, Inches(6.5), Inches(10), Inches(0.4), text, Pt(11), SUBTEXT)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# 1 封面
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, M, Inches(2.2), Inches(10), Inches(1), "QueryForge", Pt(52), WHITE, True)
txt(s, M, Inches(3.3), Inches(10), Inches(0.6), "让电商经营者获得专业数据分析能力", Pt(20), SUBTEXT)
txt(s, M, Inches(4.2), Inches(10), Inches(0.5), "ClawHunt Builder Camp 2026 · Track C: Business on AI", Pt(14), SUBTEXT)
txt(s, M, Inches(4.8), Inches(10), Inches(0.4), "面向跨境电商小团队的本地优先经营工具包 · 首个工具：订单分析 workbench", Pt(13), ACCENT)
footer(s, "queryforge-production-8d6f.up.railway.app")

# 2 痛点
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, "小卖家也需要专业分析，但没有数据团队")
bullets(s, [
    '店主问"哪个市场、品类和支付方式在拉动增长" → 手工导表',
    '运营问"哪个商品值得继续投放" → 临时筛选订单',
    '老板问"本月下滑是品类、物流还是渠道问题" → 靠经验判断',
    '广交会卖家、Amazon/marketplace 卖家和外贸 SOHO 通常不会采购 BI',
], Inches(2.0))
txt(s, M, Inches(4.8), Inches(10), Inches(0.5), "真正的痛点不是没有数据，而是经营者缺少把订单数据变成判断的工具", Pt(14), ACCENT)
footer(s, "痛点：跨境电商小团队需要专业分析能力，但预算、时间和技术门槛都很低")

# 3 方案
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, "把电商分析流程做成本地工作台")
txt(s, M, Inches(1.8), Inches(10), Inches(0.5), "自然语言提问 → 受控 SQL → 图表 → 经营建议", Pt(15), ACCENT)
bullets(s, [
    "经营者直接问订单、品类、市场、渠道和复购问题",
    "系统生成 SQL、执行只读查询，并返回图表和解释",
    "常用问题沉淀为指标，不需要每次从零分析",
    "当前 demo 是 macOS 本地包；产品形态可扩展到 Windows/macOS 预编译软件包",
], Inches(2.5))
footer(s, "不是做大型 BI，而是把专业分析能力降到小卖家也能使用的成本")

# 4 产品能力
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, "从静态报表到追问式分析")
card(s, M, Inches(1.8), Inches(5.5), Inches(2.2), "描述发生了什么",
     '"各市场本月营收是多少"\n生成图表和明细\n按市场、品类、支付方式切换')
card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(2.2), "解释为什么变化",
     '"为什么某个市场冲高后回落"\n比较时间、市场、品类、支付\n输出可验证的业务假设')
card(s, M, Inches(4.3), Inches(5.5), Inches(2.2), "验证查询可靠性",
     "SQL 先过 AST 安全校验\n只读执行，失败自动修正\n让错误暴露在流程里")
card(s, Inches(6.8), Inches(4.3), Inches(5.5), Inches(2.2), "沉淀经营指标",
     "常用问题变成指标库\n复用同一套查询口径\n形成可追问的经营记录")

# 5 真实数据验证
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, "公开 demo：Olist 巴西电商真实数据")
big(s, M, Inches(1.8), "99,441", "真实订单")
big(s, Inches(4.5), Inches(1.8), "96,096", "真实用户")
big(s, Inches(8), Inches(1.8), "74", "商品品类")
txt(s, M, Inches(3.5), Inches(10), Inches(0.4), "来自 Kaggle 公开数据集，巴西最大电商平台真实交易数据：", Pt(16), ACCENT, True)
bullets(s, [
    "总营收 R$1,601万 · 客单价 R$161 · 完成率 97% · 复购率 3.1%",
    "覆盖巴西公开交易样本，2016-2018年完整时间序列",
    "数据包含订单、支付、商品、客户地区和品类结构",
], Inches(4.0))
footer(s, "Railway / iPhone QR 是 polished public demo；完整产品路径是本地桌面软件包")

# 6 技术架构 — 三层设计
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, "桌面产品：把经营问题变成可信查询")
card(s, M, Inches(1.8), Inches(3.7), Inches(4.5), "电商指标层",
     "经营语言 → 指标定义 → SQL\n不靠 LLM 临时猜表名\n订单、品类、渠道口径可复用\n避免每次重新解释数据")
card(s, Inches(4.9), Inches(1.8), Inches(3.7), Inches(4.5), "验证式 Agent 循环",
     "理解问题 → 生成 SQL → AST 验证\n只读执行 → 读结果 → 可视化\n自纠正：出错自动修正重试\n把取数变成可交互分析")
card(s, Inches(8.9), Inches(1.8), Inches(3.7), Inches(4.5), "本地执行边界",
     "本地 DB + 只读连接\n用户自选模型服务与 key\nAST 级 SQL 安全校验\ntoken budget 留在本机")
footer(s, "必要性：经营问题需要解释、验证和追问；单纯 SQL 生成无法承担分析责任")

# 7 数据与模型边界
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, "数据与模型边界 — 本地工具也要可信")
card(s, M, Inches(1.75), Inches(5.5), Inches(2.2), "数据边界",
     "默认只暴露 schema\n只把当前问题所需查询结果交给模型解释\nAST 校验只允许 SELECT\n自动 LIMIT，避免全表扫描")
card(s, Inches(6.8), Inches(1.75), Inches(5.5), Inches(2.2), "BYOK 边界",
     "用户自带 provider key/token\nQueryForge 不转售模型能力\n公开 QR demo 不收集访客 key\n真实调用只在本地桌面端发生")
card(s, M, Inches(4.25), Inches(5.5), Inches(2.0), "商用输出依据",
     "OpenAI：用户保留 Input 权利并拥有 Output\nAnthropic：商业条款说明客户保留输出权利\nDeepSeek / Kimi：强调用户对输入、输出和授权负责")
card(s, Inches(6.8), Inches(4.25), Inches(5.5), Inches(2.0), "使用者责任",
     "老板分析自己的经营数据通常边界清楚\n个人信息、平台限制字段、第三方素材和商业秘密仍需确认权利\n输出建议不是法律/财务意见")
footer(s, "合规叙述：接入用户自己的模型账户和经营数据，不替用户绕过 provider ToS、平台条款或知识产权限制")

# 8 商业化路径
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, "商业化路径：服务小微电商的软件包")
card(s, M, Inches(1.8), Inches(5.5), Inches(2), "获客场景",
     "广交会周边地推 + 卖家社群投放\n亚马逊卖家、小工贸公司和外贸 SOHO\n现场扫码演示，几分钟看懂价值\n巴西 Olist 数据正好是海外市场样例")
card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(2), "付费方式",
     "官网/渠道/应用商店售卖闭源预编译包\n建议首发价：¥128-168\n低于一次咨询或报表外包成本\n用户自带模型 API key，token 自付")
card(s, M, Inches(4.1), Inches(5.5), Inches(2), "产品形态",
     "工具包首个工具：订单经营分析\n当前 demo 为 macOS，本质不限定平台\nOlist 是内置 demo case\n真实使用时连接本机可访问数据")
card(s, Inches(6.8), Inches(4.1), Inches(5.5), Inches(2), "开发者机会",
     "卖点不是模型本身，而是电商分析流程\n无需后期云运维，边际成本低\n低决策成本，适合冲动购买\n用小额工具费换回高频时间成本")

# 9 双重价值
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, "商业价值：让小卖家也能问经营数据")
card(s, M, Inches(1.8), Inches(5.5), Inches(4), "跨境电商小团队",
     "不用先搭 BI 或雇专职分析师\n围绕订单、品类、渠道连续追问\n快速判断海外市场和选品机会\n把常用问题沉淀成可复用指标")
card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4), "个人开发者",
     "不出售模型能力，只出售好用流程\n用户自带 provider key，成本边界清楚\n买断制降低试用门槛\n广交会等线下场景便于直接触达用户")
footer(s, "Olist 巴西电商数据用于公开 demo；真实价值是把小卖家已有经营数据变成可追问的工具包")

# 10 结束
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, M, Inches(2.2), Inches(10), Inches(1), "QueryForge", Pt(52), WHITE, True)
txt(s, M, Inches(3.3), Inches(8.4), Inches(0.8), "让电商经营者获得专业数据分析能力\n把订单、品类、渠道和复购问题变成可信查询", Pt(18), SUBTEXT)
txt(s, M, Inches(4.8), Inches(8.2), Inches(0.4), "https://queryforge-production-8d6f.up.railway.app/", Pt(14), ACCENT)
txt(s, M, Inches(5.3), Inches(8.2), Inches(0.4), "github.com/eric-stone-plus/QueryForge", Pt(12), SUBTEXT)
qr_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "assets", "qr-railway.png")
s.shapes.add_picture(qr_path, Inches(10.25), Inches(3.35), width=Inches(1.65), height=Inches(1.65))
txt(s, Inches(10.1), Inches(5.12), Inches(1.95), Inches(0.35), "扫码体验 public demo", Pt(11), SUBTEXT, a=PP_ALIGN.CENTER)
footer(s, "ClawHunt Builder Camp 2026 · Track C: Business on AI")

prs.save(os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "assets", "QueryForge-Pitch.pptx"))
print(f"Done: {len(prs.slides)} slides")
