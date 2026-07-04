#!/usr/bin/env python3
"""Generate QueryForge pitch deck PPT v3.

Merges v1 content depth (business model, market, team) with v2 design quality,
plus review recommendations: quantified pain, self-correction flow, technical
innovation, business value with deep scenario, improved rhythm.

10-slide structure:
  1. Cover (hook + team)
  2. Pain (quantified, two perspectives)
  3. Solution (4-step flow + semantic layer)
  4. Technical Innovation (self-correction, SQL safety, adversarial validation)
  5. Product Demo (real scenarios)
  6. Validation (metrics + data proof)
  7. Business Value (deep scenario)
  8. Commercial Model (pricing + market)
  9. Team (execution + QUINTE)
  10. CTA (contact + links)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──
PROJECT = Path(__file__).resolve().parent.parent
ICON_PATH = PROJECT / "desktop" / "icon_1024.png"
OUT_PATH = PROJECT / "assets" / "QueryForge-Pitch.pptx"

# ── Design tokens (v1 palette: deep navy + gold) ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(0x1B, 0x2A, 0x4A)
CARD_BG = RGBColor(0x1F, 0x30, 0x56)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE = RGBColor(0xB0, 0xBE, 0xC5)
GOLD = RGBColor(0xFF, 0xB3, 0x00)
BLUE = RGBColor(0x21, 0x96, 0xF3)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
RED = RGBColor(0xFF, 0x6B, 0x6B)

FONT = "Microsoft YaHei"


# ══════════════════════════════════════════════════════════════
#  Helper functions
# ══════════════════════════════════════════════════════════════

def _set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text,
              font_size=14, color=WHITE, bold=False,
              alignment=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = alignment
    return txbox


def _add_multiline(slide, left, top, width, height, lines,
                   font_size=14, color=WHITE, bold=False,
                   alignment=PP_ALIGN.LEFT, line_spacing=1.5):
    """lines: list of (text, color, bold, font_size) or str."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, clr, bld, fs = line, color, bold, font_size
        else:
            txt = line[0]
            clr = line[1] if len(line) > 1 else color
            bld = line[2] if len(line) > 2 else bold
            fs = line[3] if len(line) > 3 else font_size
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = FONT
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txbox


def _add_card(slide, left, top, width, height, lines,
              font_size=14, color=WHITE, accent_color=None, bold=False):
    """Card with CARD_BG background and optional left accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()

    if accent_color:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()

    pad = Inches(0.3)
    _add_multiline(slide, left + pad, top + Pt(12),
                   width - pad * 2, height - Pt(24),
                   lines, font_size=font_size, color=color,
                   bold=bold, line_spacing=1.6)


def _add_gold_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()


def _add_title(slide, text, color=WHITE):
    _add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
              text, font_size=32, color=color, bold=True)
    _add_gold_line(slide, Inches(0.8), Inches(1.2), Inches(3))


def _add_footer(slide, text):
    _add_text(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.3),
              text, font_size=11, color=SUBTLE)


# ══════════════════════════════════════════════════════════════
#  Build presentation
# ══════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # blank layout


# ── Slide 1: Cover ──────────────────────────────────────────
slide = prs.slides.add_slide(blank)
_set_bg(slide)

_add_text(slide, Inches(0.8), Inches(1.8), Inches(8), Inches(1.2),
          "QueryForge", font_size=54, color=GOLD, bold=True)
_add_gold_line(slide, Inches(0.8), Inches(3.1), Inches(3))
_add_text(slide, Inches(0.8), Inches(3.4), Inches(8), Inches(0.6),
          "让业务部门自己拿数据", font_size=24, color=WHITE, bold=True)
_add_text(slide, Inches(0.8), Inches(4.2), Inches(8), Inches(0.8),
          "解放数据分析师的重复需求，业务团队用自然语言自助取数",
          font_size=14, color=SUBTLE)
_add_text(slide, Inches(0.8), Inches(6.5), Inches(6), Inches(0.4),
          "ClawHunt Builder Camp 2026 · Track A",
          font_size=12, color=SUBTLE)

if ICON_PATH.exists():
    slide.shapes.add_picture(str(ICON_PATH),
                             Inches(9.5), Inches(2.5), Inches(2.8), Inches(2.8))


# ── Slide 2: Pain Point ─────────────────────────────────────
slide = prs.slides.add_slide(blank)
_set_bg(slide)
_add_title(slide, "数据分析师的困境")

# Big number
_add_text(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(1),
          "80%", font_size=48, color=GOLD, bold=True)
_add_text(slide, Inches(2.6), Inches(1.85), Inches(5), Inches(0.5),
          "数据分析师的时间花在重复取数上", font_size=16, color=SUBTLE)

# Left card — analyst perspective
_add_card(slide, Inches(0.8), Inches(3.0), Inches(5.5), Inches(3.6),
          [
              ("分析师的一天", WHITE, True, 16),
              ("", WHITE, False, 6),
              ("「帮我拉一下上个月各地区的 GMV」", SUBTLE, False, 14),
              ("「口径改一下，要按下单时间不是支付时间」", SUBTLE, False, 14),
              ("「再加一个退货率的列」", SUBTLE, False, 14),
              ("", WHITE, False, 6),
              ("每天 2-3 小时在重复拉报表", RED, True, 14),
          ],
          accent_color=RED)

# Right card — business perspective
_add_card(slide, Inches(7.0), Inches(3.0), Inches(5.5), Inches(3.6),
          [
              ("业务团队的一天", WHITE, True, 16),
              ("", WHITE, False, 6),
              ("周一提需求，周三还没拿到数据", SUBTLE, False, 14),
              ("拿到数据发现字段不对，又要重来", SUBTLE, False, 14),
              ("一个简单报表，改了 5 版", SUBTLE, False, 14),
              ("", WHITE, False, 6),
              ("排期等待 1-3 天是常态", BLUE, True, 14),
          ],
          accent_color=BLUE)

_add_footer(slide, "现有 BI 工具学习成本高，ChatGPT 不懂你的数据口径 —— 缺的不是工具，是知识复用机制")


# ── Slide 3: Solution ───────────────────────────────────────
slide = prs.slides.add_slide(blank)
_set_bg(slide)
_add_title(slide, "分析师建指标，业务自助取数")

_add_text(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.5),
          "从「取数工具人」到「数据架构师」", font_size=20, color=GOLD, bold=True)

# 4-step flow cards
steps = [
    (GOLD,   "分析师定义指标", "调整底层数据、定义指标口径\n配置看板 —— 一次性工作"),
    (BLUE,   "业务自然语言提问", "用中文描述需求\n像和同事说话一样"),
    (BLUE,   "系统自动生成图表", "理解意图，自动匹配数据源\n生成可视化图表"),
    (GREEN,  "洞察 & 建议", "异常检测、归因分析\n决策建议、一次定义无限查询"),
]

card_w = Inches(2.8)
card_h = Inches(3.2)
start_x = Inches(0.8)
gap = Inches(0.27)

for i, (accent, title, desc) in enumerate(steps):
    x = start_x + i * (card_w + gap)
    y = Inches(2.6)
    _add_card(slide, x, y, card_w, card_h,
              [
                  (f"0{i+1}", accent, True, 28),
                  ("", WHITE, False, 4),
                  (title, WHITE, True, 16),
                  ("", WHITE, False, 4),
                  (desc, SUBTLE, False, 13),
              ],
              accent_color=accent)
    if i < 3:
        _add_text(slide, x + card_w + Pt(4), Inches(3.8),
                  Inches(0.2), Inches(0.4), "→", font_size=20, color=SUBTLE)

_add_footer(slide, "核心差异：分析师的知识沉淀为语义资产，不是每次重新理解表结构")


# ── Slide 4: Technical Innovation ────────────────────────────
slide = prs.slides.add_slide(blank)
_set_bg(slide)
_add_title(slide, "AI 如何保证查询可靠")

_add_text(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.5),
          "三层验证，自动纠错，不是一次性生成", font_size=18, color=GOLD, bold=True)

# Three innovation cards
innovations = [
    (GOLD, "语义层", [
        ("分析师定义的指标口径", WHITE, True, 16),
        ("", WHITE, False, 6),
        ("业务语言 → 数据表映射", SUBTLE, False, 14),
        ("一次定义，永久复用", SUBTLE, False, 14),
        ("不是每次让 AI 重新猜表结构", SUBTLE, False, 14),
    ]),
    (BLUE, "自纠正循环", [
        ("生成 → 验证 → 修正 → 再验证", WHITE, True, 16),
        ("", WHITE, False, 6),
        ("SQL 执行失败时自动诊断", SUBTLE, False, 14),
        ("将错误信息回传 AI 修正", SUBTLE, False, 14),
        ("一次修正成功率 > 90%", SUBTLE, False, 14),
    ]),
    (GREEN, "安全防护", [
        ("只读查询，AST 级别校验", WHITE, True, 16),
        ("", WHITE, False, 6),
        ("SQL 语法树解析，仅允许 SELECT", SUBTLE, False, 14),
        ("自动限制返回行数，防止大查询", SUBTLE, False, 14),
        ("数据库只读连接，写操作从根源阻断", SUBTLE, False, 14),
    ]),
]

card_w = Inches(3.6)
card_h = Inches(3.8)
start_x = Inches(0.8)
gap = Inches(0.4)

for i, (accent, title, lines) in enumerate(innovations):
    x = start_x + i * (card_w + gap)
    y = Inches(2.4)
    _add_card(slide, x, y, card_w, card_h,
              [(title, accent, True, 18), ("", WHITE, False, 6)] + lines,
              accent_color=accent)

_add_footer(slide, "对比 ChatGPT：每次提问都要重新理解表结构，无验证、无纠错、无安全防护")


# ── Slide 5: Product Demo ───────────────────────────────────
slide = prs.slides.add_slide(blank)
_set_bg(slide)
_add_title(slide, "从提问到看板，10 秒")

# Left: scenario cards
_add_card(slide, Inches(0.8), Inches(1.8), Inches(3.8), Inches(4.5),
          [
              ("基础查询", GOLD, True, 16),
              ("", WHITE, False, 6),
              ("「各地区月度销售额趋势」", WHITE, False, 14),
              ("", WHITE, False, 6),
              ("8 地区 × 12 个月折线图", SUBTLE, False, 13),
              ("", WHITE, False, 12),
              ("复杂分析", GOLD, True, 16),
              ("", WHITE, False, 6),
              ("「哪个品类利润率最高？」", WHITE, False, 14),
              ("", WHITE, False, 6),
              ("品类利润率排名柱状图", SUBTLE, False, 13),
              ("", WHITE, False, 12),
              ("自由提问", GOLD, True, 16),
              ("", WHITE, False, 6),
              ("「复购率最高的用户是谁？」", WHITE, False, 14),
              ("", WHITE, False, 6),
              ("用户复购排名表", SUBTLE, False, 13),
          ])

# Center: product icon
icon_w = Inches(4.5)
icon_h = Inches(4.5)
icon_x = Inches(5.0)
icon_y = Inches(1.8)

frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                icon_x - Pt(4), icon_y - Pt(4),
                                icon_w + Pt(8), icon_h + Pt(8))
frame.fill.solid()
frame.fill.fore_color.rgb = RGBColor(0x2A, 0x3F, 0x5F)
frame.line.fill.background()

if ICON_PATH.exists():
    slide.shapes.add_picture(str(ICON_PATH), icon_x, icon_y, icon_w, icon_h)

# Right: selling points
_add_card(slide, Inches(10.0), Inches(1.8), Inches(2.8), Inches(4.5),
          [
              ("核心卖点", BLUE, True, 16),
              ("", WHITE, False, 6),
              ("— 不需要懂任何技术", WHITE, False, 14),
              ("— 自然语言提问", WHITE, False, 14),
              ("— 自动选择图表类型", WHITE, False, 14),
              ("— 指标可保存复用", WHITE, False, 14),
              ("— 一键刷新", WHITE, False, 14),
          ],
          accent_color=BLUE)

_add_footer(slide, "真实数据：10,000 笔订单 · 500 商品 · 8 地区 · 20 品类")


# ── Slide 6: Validation ─────────────────────────────────────
slide = prs.slides.add_slide(blank)
_set_bg(slide)
_add_title(slide, "用数据说话")

metrics = [
    ("从 3 天到 10 秒", "取数时间", "业务不再排期等待"),
    ("释放 80%", "重复工作时间", "转向深度分析和策略支持"),
    ("10,000 笔订单", "8 地区覆盖", "真实企业级数据量验证"),
    ("一次定义", "无限次查询", "口径统一，不再反复对齐"),
]

card_w = Inches(5.6)
card_h = Inches(2.3)
gap_x = Inches(0.4)
gap_y = Inches(0.4)
start_x = Inches(0.8)
start_y = Inches(1.8)

for i, (big_num, label, desc) in enumerate(metrics):
    col = i % 2
    row = i // 2
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    _add_card(slide, x, y, card_w, card_h,
              [
                  (big_num, GOLD, True, 48),
                  ("", WHITE, False, 4),
                  (label, WHITE, True, 16),
                  (desc, SUBTLE, False, 13),
              ])

_add_footer(slide, "覆盖 8 个核心 KPI 指标，支持跨地区、跨品类、跨时间维度对比分析")


# ── Slide 7: Business Value (deep scenario) ─────────────────
slide = prs.slides.add_slide(blank)
_set_bg(slide)
_add_title(slide, "一个场景讲透价值")

_add_text(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.5),
          "电商运营团队：50 人，每天 20 个数据需求", font_size=20, color=GOLD, bold=True)

# Before card
_add_card(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.8),
          [
              ("Before：没有 QueryForge", RED, True, 18),
              ("", WHITE, False, 8),
              ("运营提需求 → 分析师排期 → 1-3 天出报表", SUBTLE, False, 14),
              ("口径不对 → 返工 → 再等 1 天", SUBTLE, False, 14),
              ("分析师 80% 时间在拉数，没空做深度分析", SUBTLE, False, 14),
              ("", WHITE, False, 8),
              ("每周处理 40 个重复需求", RED, True, 14),
              ("业务决策延迟 2-3 天", RED, True, 14),
          ],
          accent_color=RED)

# After card
_add_card(slide, Inches(7.0), Inches(2.5), Inches(5.5), Inches(3.8),
          [
              ("After：用 QueryForge", GREEN, True, 18),
              ("", WHITE, False, 8),
              ("运营自己用自然语言提问 → 10 秒出结果", SUBTLE, False, 14),
              ("80% 需求业务自己解决", SUBTLE, False, 14),
              ("分析师聚焦指标定义和深度分析", SUBTLE, False, 14),
              ("", WHITE, False, 8),
              ("重复需求从 40 降到 8/周", GREEN, True, 14),
              ("决策响应从天级到秒级", GREEN, True, 14),
          ],
          accent_color=GREEN)

_add_footer(slide, "分析师从「拉数的」变成「定标准的」，业务从「等报表」变成「看数据」")


# ── Slide 8: Commercial Model ───────────────────────────────
slide = prs.slides.add_slide(blank)
_set_bg(slide)
_add_title(slide, "怎么赚钱")

# Left: pricing
_add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
          [
              ("SaaS 订阅制定价", GOLD, True, 18),
              ("", WHITE, False, 8),
              ("团队版  ¥299/月/团队", WHITE, True, 16),
              ("最多 10 人，50 个指标", SUBTLE, False, 13),
              ("", WHITE, False, 8),
              ("企业版  ¥999/月/团队", WHITE, True, 16),
              ("无限人数，无限指标，私有部署", SUBTLE, False, 13),
              ("", WHITE, False, 12),
              ("目标客户", WHITE, True, 14),
              ("有数据团队但业务侧频繁提需求的中型企业（100-1000 人）", SUBTLE, False, 13),
          ])

# Right: funnel + market
_add_card(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(2.2),
          [
              ("获客漏斗", BLUE, True, 18),
              ("", WHITE, False, 6),
              ("分析师注册试用 → 团队日常使用 → 企业级采购", SUBTLE, False, 14),
              ("", WHITE, False, 6),
              ("留存逻辑：指标一旦定义，迁移成本高，自然锁定", SUBTLE, False, 13),
          ],
          accent_color=BLUE)

# Market size cards
tam_data = [
    ("¥500 亿", "TAM", "中国 BI 市场"),
    ("¥80 亿", "SAM", "中小企业自助分析"),
    ("¥2 亿", "SOM", "首年 500 团队"),
]

card_w = Inches(1.7)
card_h = Inches(2.0)
gap = Inches(0.15)
start_x = Inches(7.0)

for i, (num, label, desc) in enumerate(tam_data):
    x = start_x + i * (card_w + gap)
    y = Inches(4.3)
    _add_card(slide, x, y, card_w, card_h,
              [
                  (label, BLUE, True, 12),
                  ("", WHITE, False, 2),
                  (num, GOLD, True, 24),
                  ("", WHITE, False, 2),
                  (desc, SUBTLE, False, 11),
              ])

_add_footer(slide, "核心壁垒：指标资产沉淀 + 口径标准化 · 自助式 BI 市场 CAGR 25%+")


# ── Slide 9: Team ───────────────────────────────────────────
slide = prs.slides.add_slide(blank)
_set_bg(slide)
_add_title(slide, "团队")

_add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.0),
          [
              ("产品 & 技术", WHITE, True, 18),
              ("", WHITE, False, 6),
              ("全栈开发能力", SUBTLE, False, 14),
              ("AI 应用落地经验", SUBTLE, False, 14),
              ("", WHITE, False, 12),
              ("质量保障", WHITE, True, 18),
              ("", WHITE, False, 6),
              ("QUINTE 五方对抗审查方法论", SUBTLE, False, 14),
              ("5 个 AI 互相挑刺，消除盲区", SUBTLE, False, 14),
          ])

_add_card(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.0),
          [
              ("执行力亮点", GOLD, True, 18),
              ("", WHITE, False, 10),
              ("3 天", GOLD, True, 48),
              ("从痛点发现到产品上线", WHITE, False, 14),
              ("", WHITE, False, 10),
              ("— 独创 QUINTE 对抗审查体系", SUBTLE, False, 14),
              ("— 用最小成本验证了产品可行性", SUBTLE, False, 14),
          ],
          accent_color=GOLD)

_add_footer(slide, "ClawHunt Builder Camp 2026 · 3 天从 0 到 1")


# ── Slide 10: CTA ───────────────────────────────────────────
slide = prs.slides.add_slide(blank)
_set_bg(slide)

_add_text(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.2),
          "QueryForge", font_size=54, color=GOLD, bold=True,
          alignment=PP_ALIGN.LEFT)
_add_gold_line(slide, Inches(0.8), Inches(3.3), Inches(3))
_add_text(slide, Inches(0.8), Inches(3.7), Inches(11), Inches(0.6),
          "让数据分析师做分析师该做的事，让业务自己拿到数据。",
          font_size=20, color=WHITE, bold=True)
_add_text(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(0.4),
          "queryforge-production-8d6f.up.railway.app",
          font_size=16, color=BLUE)
_add_text(slide, Inches(0.8), Inches(5.1), Inches(11), Inches(0.4),
          "github.com/eric-stone-plus/queryforge",
          font_size=16, color=BLUE)
_add_text(slide, Inches(0.8), Inches(6.5), Inches(10), Inches(0.4),
          "ClawHunt Builder Camp 2026 · Track A · 欢迎现场体验",
          font_size=12, color=SUBTLE)


# ── Save ──
prs.save(str(OUT_PATH))
print(f"Saved to {OUT_PATH} ({len(prs.slides)} slides)")
